#!/usr/bin/env python3
"""Convert a slide image into an editable-text PPTX using a layout JSON.

This script implements the high-fidelity method used by the skill:
1. full-slide original image as background
2. optional cover patches over original text regions
3. visible editable textboxes over the patches

Layout coordinates are in source image pixels.
"""
import argparse
import json
import os
from pathlib import Path
from typing import Any, Dict

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


def hex_to_rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def parse_align(value: str):
    return {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get((value or "left").lower(), PP_ALIGN.LEFT)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("image", help="Source slide image, e.g. input.png")
    parser.add_argument("output", help="Output PPTX path")
    parser.add_argument("--layout", required=True, help="Layout JSON with patches and textboxes")
    args = parser.parse_args()

    image_path = Path(args.image)
    output_path = Path(args.output)
    layout: Dict[str, Any] = json.loads(Path(args.layout).read_text(encoding="utf-8"))

    img = Image.open(image_path)
    img_w, img_h = img.size

    prs = Presentation()
    prs.slide_width = Inches(float(layout.get("slide_width_in", 13.333333)))
    prs.slide_height = Inches(float(layout.get("slide_height_in", 7.5)))
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)

    sx = prs.slide_width.inches / img_w
    sy = prs.slide_height.inches / img_h

    def X(px): return Inches(float(px) * sx)
    def Y(py): return Inches(float(py) * sy)
    def Wd(px): return Inches(float(px) * sx)
    def Hd(py): return Inches(float(py) * sy)

    # Add cover patches. Use shapes here; if renderer adds shadows, replace with a white PNG patch.
    for patch in layout.get("patches", []):
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            X(patch["x"]), Y(patch["y"]), Wd(patch["w"]), Hd(patch["h"])
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(patch.get("color", "#FFFFFF"))
        shape.line.fill.background()

    # Add visible editable textboxes.
    for item in layout.get("texts", []):
        tx = slide.shapes.add_textbox(X(item["x"]), Y(item["y"]), Wd(item["w"]), Hd(item["h"]))
        tx.fill.background()
        tx.line.fill.background()
        tf = tx.text_frame
        tf.clear()
        tf.word_wrap = bool(item.get("word_wrap", True))
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        paragraphs = item.get("paragraphs")
        if paragraphs is None:
            paragraphs = [{"runs": item.get("runs", []), "align": item.get("align", "left")}]

        for p_idx, p_spec in enumerate(paragraphs):
            p = tf.paragraphs[0] if p_idx == 0 else tf.add_paragraph()
            p.alignment = parse_align(p_spec.get("align", item.get("align", "left")))
            for run_spec in p_spec.get("runs", []):
                r = p.add_run()
                r.text = run_spec.get("text", "")
                r.font.name = run_spec.get("font", item.get("font", "Microsoft YaHei"))
                r.font.size = Pt(float(run_spec.get("size", item.get("size", 18))))
                r.font.bold = bool(run_spec.get("bold", item.get("bold", False)))
                r.font.color.rgb = hex_to_rgb(run_spec.get("color", item.get("color", "#000000")))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(output_path)


if __name__ == "__main__":
    main()
