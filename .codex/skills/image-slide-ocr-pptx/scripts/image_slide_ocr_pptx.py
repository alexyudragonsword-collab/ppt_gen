#!/usr/bin/env python3
"""
image_slide_ocr_pptx.py

Convert image-only slide sources into an editable-text PPTX while preserving the
original slide image as a full-page background.

Supported inputs:
  - one or more image files: png, jpg, jpeg, webp
  - a directory containing image files
  - an image-only pptx where each slide contains a full-slide or dominant picture
  - optional OCR JSON produced by a vision model or another OCR pipeline

Examples:
  python image_slide_ocr_pptx.py input_dir --output editable.pptx
  python image_slide_ocr_pptx.py image1.png image2.png --output editable.pptx
  python image_slide_ocr_pptx.py image_only.pptx --output editable.pptx --lang chi_sim+eng
  python image_slide_ocr_pptx.py slides/ --ocr-json ocr_blocks.json --output editable.pptx

OCR JSON format:
[
  [
    {"text": "Title", "x": 0.08, "y": 0.06, "w": 0.7, "h": 0.08,
     "font_size": 28, "bold": true, "color": "#111111", "align": "left"}
  ],
  [ ... page 2 blocks ... ]
]

Coordinate fields x/y/w/h are normalized to page width/height, from 0 to 1.
"""

from __future__ import annotations

import argparse
import json
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable, List, Optional, Sequence, Tuple

from PIL import Image, ImageStat
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
DEFAULT_SLIDE_WIDTH_IN = 13.333333


@dataclass
class TextBlock:
    text: str
    x: float
    y: float
    w: float
    h: float
    font_size: float = 14
    bold: bool = False
    color: str = "#000000"
    align: str = "left"
    font: str = "Microsoft YaHei"

    @classmethod
    def from_dict(cls, data: dict[str, Any]) -> "TextBlock":
        return cls(
            text=str(data.get("text", "")).strip(),
            x=float(data.get("x", 0)),
            y=float(data.get("y", 0)),
            w=float(data.get("w", 0)),
            h=float(data.get("h", 0)),
            font_size=float(data.get("font_size", 14)),
            bold=bool(data.get("bold", False)),
            color=str(data.get("color", "#000000")),
            align=str(data.get("align", "left")),
            font=str(data.get("font", "Microsoft YaHei")),
        )

    def to_dict(self) -> dict[str, Any]:
        return {
            "text": self.text,
            "x": self.x,
            "y": self.y,
            "w": self.w,
            "h": self.h,
            "font_size": self.font_size,
            "bold": self.bold,
            "color": self.color,
            "align": self.align,
            "font": self.font,
        }


def natural_key(path: Path) -> list[Any]:
    parts = re.split(r"(\d+)", path.name.lower())
    return [int(p) if p.isdigit() else p for p in parts]


def clamp(value: float, lo: float, hi: float) -> float:
    return max(lo, min(hi, value))


def sanitize_hex_color(value: str, fallback: str = "#000000") -> str:
    value = (value or "").strip()
    if re.fullmatch(r"#[0-9a-fA-F]{6}", value):
        return value
    if re.fullmatch(r"[0-9a-fA-F]{6}", value):
        return "#" + value
    return fallback


def hex_to_rgb(value: str) -> tuple[int, int, int]:
    value = sanitize_hex_color(value).lstrip("#")
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def is_image_file(path: Path) -> bool:
    return path.is_file() and path.suffix.lower() in IMAGE_EXTS


def collect_image_inputs(inputs: Sequence[Path]) -> list[Path]:
    images: list[Path] = []
    for item in inputs:
        if item.is_dir():
            for child in sorted(item.iterdir(), key=natural_key):
                if is_image_file(child):
                    images.append(child)
        elif is_image_file(item):
            images.append(item)
    return sorted(images, key=natural_key)


def extract_images_from_image_only_pptx(pptx_path: Path, workdir: Path) -> list[Path]:
    """Extract dominant/full-slide images from a PPTX using python-pptx.

    This works well for PPTX files where each slide is one large picture. It does
    not render arbitrary vector/text slides. For those, export the PPTX to images
    first, then run this script on the exported images.
    """
    prs = Presentation(str(pptx_path))
    out_dir = workdir / f"{pptx_path.stem}_slides"
    out_dir.mkdir(parents=True, exist_ok=True)

    extracted: list[Path] = []
    for idx, slide in enumerate(prs.slides, start=1):
        candidates = []
        for shape in slide.shapes:
            if not hasattr(shape, "image"):
                continue
            area = int(shape.width) * int(shape.height)
            candidates.append((area, shape))
        if not candidates:
            continue
        _, shape = max(candidates, key=lambda item: item[0])
        blob = shape.image.blob
        ext = shape.image.ext or "png"
        image_path = out_dir / f"slide_{idx:03d}.{ext}"
        image_path.write_bytes(blob)
        extracted.append(image_path)
    return extracted


def try_render_pptx_with_libreoffice(pptx_path: Path, workdir: Path) -> list[Path]:
    """Best-effort PPTX rendering with LibreOffice if available.

    This fallback is useful when the PPTX has shapes rather than a single image.
    Some environments do not have LibreOffice, so this function is optional.
    """
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return []

    pdf_dir = workdir / f"{pptx_path.stem}_pdf"
    img_dir = workdir / f"{pptx_path.stem}_rendered"
    pdf_dir.mkdir(parents=True, exist_ok=True)
    img_dir.mkdir(parents=True, exist_ok=True)

    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(pdf_dir), str(pptx_path)],
        check=False,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
    )
    pdfs = list(pdf_dir.glob("*.pdf"))
    if not pdfs:
        return []

    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        return []

    prefix = img_dir / "slide"
    subprocess.run(
        [pdftoppm, "-png", "-r", "220", str(pdfs[0]), str(prefix)],
        check=False,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
    )
    return sorted(img_dir.glob("slide-*.png"), key=natural_key)


def resolve_inputs(inputs: Sequence[str], workdir: Path) -> list[Path]:
    paths = [Path(p).expanduser().resolve() for p in inputs]
    missing = [str(p) for p in paths if not p.exists()]
    if missing:
        raise FileNotFoundError("Input path not found: " + ", ".join(missing))

    images: list[Path] = []
    pptx_paths: list[Path] = []
    for path in paths:
        if path.suffix.lower() == ".pptx":
            pptx_paths.append(path)
        elif path.is_dir() or is_image_file(path):
            images.extend(collect_image_inputs([path]))
        else:
            raise ValueError(f"Unsupported input type: {path}")

    for pptx_path in pptx_paths:
        extracted = extract_images_from_image_only_pptx(pptx_path, workdir)
        if not extracted:
            extracted = try_render_pptx_with_libreoffice(pptx_path, workdir)
        if not extracted:
            raise RuntimeError(
                f"Could not extract or render slides from {pptx_path}. "
                "Export it to images first, then rerun this script."
            )
        images.extend(extracted)

    if not images:
        raise RuntimeError("No image slides were found in the provided inputs.")
    return sorted(images, key=natural_key)


def estimate_text_color(image: Image.Image, box: tuple[int, int, int, int]) -> str:
    """Choose black or white by local background brightness.

    This is a practical approximation for overlay text. Users can supply OCR JSON
    with exact colors when visual fidelity matters.
    """
    x, y, w, h = box
    pad = max(2, int(min(w, h) * 0.15))
    crop_box = (
        max(0, x - pad),
        max(0, y - pad),
        min(image.width, x + w + pad),
        min(image.height, y + h + pad),
    )
    crop = image.convert("L").crop(crop_box)
    mean = ImageStat.Stat(crop).mean[0] if crop.size[0] and crop.size[1] else 255
    return "#000000" if mean > 150 else "#FFFFFF"


def group_tesseract_words_into_lines(data: dict[str, list[Any]], image: Image.Image, min_conf: int) -> list[TextBlock]:
    rows: dict[tuple[int, int, int], list[int]] = {}
    n = len(data.get("text", []))

    for i in range(n):
        text = str(data["text"][i]).strip()
        if not text:
            continue
        try:
            conf = float(data.get("conf", ["-1"])[i])
        except Exception:
            conf = -1
        if conf < min_conf:
            continue
        key = (
            int(data.get("block_num", [0])[i]),
            int(data.get("par_num", [0])[i]),
            int(data.get("line_num", [0])[i]),
        )
        rows.setdefault(key, []).append(i)

    blocks: list[TextBlock] = []
    for _, indices in sorted(rows.items()):
        words = [str(data["text"][i]).strip() for i in indices if str(data["text"][i]).strip()]
        if not words:
            continue
        text = " ".join(words)

        left = min(int(data["left"][i]) for i in indices)
        top = min(int(data["top"][i]) for i in indices)
        right = max(int(data["left"][i]) + int(data["width"][i]) for i in indices)
        bottom = max(int(data["top"][i]) + int(data["height"][i]) for i in indices)
        width = max(1, right - left)
        height = max(1, bottom - top)

        # Add a small margin so text is not clipped in PowerPoint.
        margin_x = int(width * 0.03) + 2
        margin_y = int(height * 0.20) + 2
        left = max(0, left - margin_x)
        top = max(0, top - margin_y)
        width = min(image.width - left, width + margin_x * 2)
        height = min(image.height - top, height + margin_y * 2)

        font_size = max(6, min(44, height * 0.55))
        color = estimate_text_color(image, (left, top, width, height))

        blocks.append(
            TextBlock(
                text=text,
                x=clamp(left / image.width, 0, 1),
                y=clamp(top / image.height, 0, 1),
                w=clamp(width / image.width, 0.001, 1),
                h=clamp(height / image.height, 0.001, 1),
                font_size=font_size,
                bold=False,
                color=color,
                align="left",
            )
        )

    return merge_nearby_short_lines(blocks)


def merge_nearby_short_lines(blocks: list[TextBlock]) -> list[TextBlock]:
    """Conservative line cleanup to reduce excessive tiny textboxes."""
    if not blocks:
        return blocks
    # Keep the implementation conservative: do not merge distant lines or columns.
    cleaned: list[TextBlock] = []
    for block in blocks:
        if not block.text.strip():
            continue
        if len(block.text) == 1 and not block.text.isalnum():
            continue
        cleaned.append(block)
    return cleaned


def run_tesseract_ocr(image_path: Path, lang: str, min_conf: int) -> list[TextBlock]:
    try:
        import pytesseract  # type: ignore
        from pytesseract import Output  # type: ignore
    except Exception as exc:
        raise RuntimeError(
            "pytesseract is not installed. Install it or provide --ocr-json. "
            "Python package: pip install pytesseract. System binary: install tesseract-ocr."
        ) from exc

    image = Image.open(image_path).convert("RGB")
    try:
        data = pytesseract.image_to_data(image, lang=lang, output_type=Output.DICT)
    except Exception as exc:
        raise RuntimeError(
            f"Tesseract OCR failed for {image_path}. "
            "Check that the tesseract binary and requested language data are installed, "
            "or provide --ocr-json."
        ) from exc
    return group_tesseract_words_into_lines(data, image, min_conf=min_conf)


def load_ocr_json(path: Path) -> list[list[TextBlock]]:
    raw = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(raw, list):
        raise ValueError("OCR JSON must be a list of pages.")
    pages: list[list[TextBlock]] = []
    for page in raw:
        if not isinstance(page, list):
            raise ValueError("Each OCR JSON page must be a list of text blocks.")
        pages.append([TextBlock.from_dict(item) for item in page if isinstance(item, dict)])
    return pages


def save_ocr_json(path: Path, pages: Sequence[Sequence[TextBlock]]) -> None:
    payload = [[block.to_dict() for block in page] for page in pages]
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def fit_slide_size_to_first_image(prs: Presentation, first_image: Path) -> None:
    with Image.open(first_image) as img:
        w, h = img.size
    prs.slide_width = Inches(DEFAULT_SLIDE_WIDTH_IN)
    prs.slide_height = Inches(DEFAULT_SLIDE_WIDTH_IN * h / w)


def add_picture_background(slide, image_path: Path, slide_w: int, slide_h: int) -> None:
    slide.shapes.add_picture(str(image_path), 0, 0, width=slide_w, height=slide_h)


def add_text_block(slide, block: TextBlock, slide_w: int, slide_h: int, default_font: str) -> None:
    text = block.text.strip()
    if not text:
        return

    x = int(clamp(block.x, 0, 1) * slide_w)
    y = int(clamp(block.y, 0, 1) * slide_h)
    w = int(clamp(block.w, 0.001, 1) * slide_w)
    h = int(clamp(block.h, 0.001, 1) * slide_h)

    textbox = slide.shapes.add_textbox(x, y, w, h)
    textbox.fill.background()
    textbox.line.fill.background()

    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(block.align.lower(), PP_ALIGN.LEFT)

    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
        run.text = text

    run.font.name = block.font or default_font
    run.font.size = Pt(max(4, block.font_size))
    run.font.bold = bool(block.bold)
    r, g, b = hex_to_rgb(block.color)
    run.font.color.rgb = RGBColor(r, g, b)


def build_pptx(
    image_paths: Sequence[Path],
    ocr_pages: Sequence[Sequence[TextBlock]],
    output_path: Path,
    default_font: str,
) -> None:
    prs = Presentation()
    fit_slide_size_to_first_image(prs, image_paths[0])

    blank = prs.slide_layouts[6]
    for image_path, blocks in zip(image_paths, ocr_pages):
        slide = prs.slides.add_slide(blank)
        add_picture_background(slide, image_path, prs.slide_width, prs.slide_height)
        for block in blocks:
            add_text_block(slide, block, prs.slide_width, prs.slide_height, default_font)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def inspect_pptx(path: Path) -> list[dict[str, int]]:
    prs = Presentation(str(path))
    rows: list[dict[str, int]] = []
    for i, slide in enumerate(prs.slides, start=1):
        text_shapes = [s for s in slide.shapes if hasattr(s, "text") and str(s.text).strip()]
        pic_shapes = [s for s in slide.shapes if getattr(s, "shape_type", None) == 13]
        rows.append({"page": i, "pictures": len(pic_shapes), "textboxes": len(text_shapes)})
    return rows


def main(argv: Optional[Sequence[str]] = None) -> int:
    parser = argparse.ArgumentParser(
        description="Rebuild image slides into PPTX with full-slide image backgrounds and editable OCR textboxes."
    )
    parser.add_argument("inputs", nargs="+", help="Image file(s), image directory, or image-only PPTX file(s).")
    parser.add_argument("-o", "--output", required=True, help="Output PPTX path.")
    parser.add_argument("--ocr-json", help="Optional OCR JSON file with per-page normalized text blocks.")
    parser.add_argument("--save-ocr-json", help="Write detected OCR blocks to this JSON file for manual correction.")
    parser.add_argument("--lang", default="chi_sim+eng", help="Tesseract language, default: chi_sim+eng.")
    parser.add_argument("--min-conf", type=int, default=35, help="Minimum Tesseract confidence, default: 35.")
    parser.add_argument("--font", default="Microsoft YaHei", help="Default font for generated textboxes.")
    parser.add_argument("--no-ocr", action="store_true", help="Only preserve slide images; do not create textboxes.")
    parser.add_argument("--keep-workdir", action="store_true", help="Do not delete temporary extracted images.")
    args = parser.parse_args(argv)

    tmp = tempfile.TemporaryDirectory(prefix="image_slide_ocr_")
    workdir = Path(tmp.name)

    try:
        image_paths = resolve_inputs(args.inputs, workdir)

        if args.ocr_json:
            ocr_pages = load_ocr_json(Path(args.ocr_json).expanduser().resolve())
            if len(ocr_pages) != len(image_paths):
                raise ValueError(
                    f"OCR JSON page count ({len(ocr_pages)}) does not match image count ({len(image_paths)})."
                )
        elif args.no_ocr:
            ocr_pages = [[] for _ in image_paths]
        else:
            ocr_pages = []
            for i, image_path in enumerate(image_paths, start=1):
                print(f"OCR page {i}/{len(image_paths)}: {image_path}", file=sys.stderr)
                try:
                    blocks = run_tesseract_ocr(image_path, lang=args.lang, min_conf=args.min_conf)
                except RuntimeError as exc:
                    print(f"WARNING: {exc}", file=sys.stderr)
                    blocks = []
                ocr_pages.append(blocks)

        if args.save_ocr_json:
            save_ocr_json(Path(args.save_ocr_json).expanduser().resolve(), ocr_pages)

        output_path = Path(args.output).expanduser().resolve()
        build_pptx(image_paths, ocr_pages, output_path, default_font=args.font)
        inspection = inspect_pptx(output_path)

        print(f"Created: {output_path}")
        print(f"Pages: {len(inspection)}")
        for row in inspection:
            print(f"Page {row['page']}: pictures={row['pictures']} textboxes={row['textboxes']}")
        return 0
    finally:
        if args.keep_workdir:
            print(f"Workdir kept: {workdir}", file=sys.stderr)
        else:
            tmp.cleanup()


if __name__ == "__main__":
    raise SystemExit(main())
